from email.message import EmailMessage
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
import time 
import msal
import requests
from powerbi.client import PowerBiClient
import img2pdf
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.support.ui import WebDriverWait
from selenium.common.exceptions import NoSuchElementException
import os
import configparser


class Authentification_for_PowerBI:
    def __init__(self,client_id,user,passwd,tenant,client_secret,file):
        self.CLIENT_ID=client_id
        self.username=user
        self.password=passwd
        self.TENANT_ID=tenant
        self.AUTHORITY_URL = 'https://login.microsoftonline.com/'+self.TENANT_ID
        self.SCOPE = ["https://analysis.windows.net/powerbi/api/.default"]
        self.URL_TO_GET_GROUPS = 'https://api.powerbi.com/v1.0/myorg/groups'
        self.file=file
        
        self.CLIENT_POWER_BI = PowerBiClient(
            client_id=self.CLIENT_ID,
            client_secret=client_secret,
            scope=['https://analysis.windows.net/powerbi/api/.default'],
            redirect_uri="https://localhost/redirect",
            credentials=self.file
        )

        self.reports_service = self.CLIENT_POWER_BI.reports()
    def get_my_workspace(self):
        return self.reports_service.get_reports()

    def get_tenant(self):
        return str(self.TENANT_ID)

    def get_workspace_informations(self):
        client_data = msal.PublicClientApplication(self.CLIENT_ID, authority=self.AUTHORITY_URL)
        result = client_data.acquire_token_by_username_password(username=self.username,password=self.password,scopes=self.SCOPE)
        if 'access_token' in result:
            access_token = result['access_token']
            header = {'Content-Type':'application/json','Authorization': f'Bearer {access_token}'}
            api_out = requests.get(url=self.URL_TO_GET_GROUPS, headers=header)
            return api_out.json()
        else:
            return result.get("error")

    def show_workspace(self):
        work_space=[]
        for i in self.get_workspace_informations()['value']:
            work_space.append([i['name'],i['id']])
        return work_space

    ##  return reports id in each workspace
    def get_report_from_workspace(self,id_work_space):
        report_out=self.reports_service.get_group_reports(
            group_id=str(id_work_space)
        )
        return report_out
                
    def grab_my_report(self,REPORT_ID):
        return self.reports_service.get_report(report_id=str(REPORT_ID))

    def all_report(self):
        try:
            l=[]
            nb=0
            id_work_space=self.show_workspace()
            for i in id_work_space:
                l.append([i[0],self.get_report_from_workspace(i[1])])
                nb+=len([_['name'] for _ in self.get_report_from_workspace(i[1])['value']])

            nb+=len([i['name'] for i in self.get_my_workspace()['value']])
            return [l,nb]
        except:
            return self.show_workspace()

import os
def screen_matidhekech(url_link,email,passwd):
    kkk=[]
    for e in range(len(url_link.split(','))):
        url_link.split(',')[e]+='/ReportSection'
        print(url_link.split(',')[e])
        print("path",os.path.abspath(os.getcwd()))
        chrome_options = Options() #
        chrome_options.add_argument('--headless')

        chrome_options.add_argument("--window-size=1920x1080")
        driver = webdriver.Chrome(options=chrome_options, executable_path=os.path.abspath(os.getcwd())+"/chromedriver.exe")
        # service object
        driver.get(url_link.split(',')[e]+'/ReportSection')
        time.sleep(6)
        driver.find_element_by_id('email').send_keys(email)
        driver.find_element_by_id('submitBtn').click()
        l=WebDriverWait(driver=driver, timeout=10).until(
            lambda x: x.execute_script("return document.readyState === 'complete'")
        )
        if l:
            try:
                time.sleep(6)
                driver.find_element_by_name('passwd').send_keys(passwd)
                driver.find_element_by_id('idSIButton9').click()
                g=WebDriverWait(driver=driver, timeout=12).until(
                    lambda x: x.execute_script("return document.readyState === 'complete'")
                )
            except:
                time.sleep(9)
                driver.find_element_by_name('passwd').send_keys(passwd)
                driver.find_element_by_id('idSIButton9').click()
                g=WebDriverWait(driver=driver, timeout=12).until(
                    lambda x: x.execute_script("return document.readyState === 'complete'")
                )
            if g:
                time.sleep(6)
                driver.find_element_by_id('idSIButton9').click()
                h=WebDriverWait(driver=driver, timeout=10).until(
                    lambda x: x.execute_script("return document.readyState === 'complete'")
                )
                if h:
                    time.sleep(8)
                    file=[]
                    try:
                        element = driver.find_element_by_class_name("displayAreaViewport")
                        try:
                            pages_data=driver.find_element_by_tag_name('mat-list')
                            items = pages_data.find_elements_by_tag_name("li")
                            i=0
                            for item in items:
                                item.click()
                                hhhhh=WebDriverWait(driver=driver, timeout=10).until(
                                    lambda x: x.execute_script("return document.readyState === 'complete'")
                                )
                                if hhhhh:
                                    time.sleep(6)
                                    i+=1
                                    path_pic=f'{os.path.abspath(os.getcwd())}\\dashbord\\static\\blog\\files\\page{i}{e}.png'
                                    element.screenshot(path_pic)
                                    file.append(path_pic)
                            return file

                        except NoSuchElementException:
                            path_pic=f'{os.path.abspath(os.getcwd())}\\dashbord\\static\\blog\\files\\page{0}{e}.png'
                            element.screenshot(path_pic)
                            kkk.append(path_pic)
                            
                    except NoSuchElementException:
                        return
    driver.close()
    return kkk

from pptx import Presentation

def transform_file_to_pdf(name_folder,pict_list,format):
    if format == "PDF File":
        pdf_name_path=f'{os.path.abspath(os.getcwd())}\\dashbord\\static\\blog\\files\\'+name_folder+".pdf"
        print(pdf_name_path)
        with open(pdf_name_path,"wb") as f:
            f.write(img2pdf.convert(pict_list))
        return pdf_name_path

    if format=="pptx":
        p = Presentation()
        ppt_name_path=os.path.abspath(os.getcwd())+"\\dashbord\\static\\blog\\files\\"+name_folder+".pptx"
        ims = pict_list
        blank_slide_layout= p.slide_layouts[6]
        for im in ims:
            slide = p.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(im, 0, 0, p.slide_width, p.slide_height)
        p.save(ppt_name_path)
        return ppt_name_path

############################################################  envoyer pbi report ou pdf a une email 
def send_pdf_file(recieve,subject,path,mesg):
    k=recieve.split(",")
    
    for i in k:
        body = f'{mesg},\npiece jointe:'
        sender,password = 'centre.declaration@gmail.com','rcrtbuvmjyxkofcj'
        message = MIMEMultipart()
        message['From'] = sender
        message['To'] = i
        message['Subject'] = subject
        message.attach(MIMEText(body, 'plain'))
        pdfname=path
        binary_pdf = open(pdfname, 'rb')
        payload = MIMEBase('application', 'octate-stream', Name=pdfname)
        payload.set_payload((binary_pdf).read())
        encoders.encode_base64(payload)
        payload.add_header('Content-Decomposition', 'attachment', filename=pdfname)
        message.attach(payload)        
        session = smtplib.SMTP('smtp.gmail.com', 587)
        session.starttls()
        session.login(sender, password)
        text = message.as_string()
        session.sendmail(sender, i, text)
        session.quit()    
    return True


def verif_msg(randnbr,email):
    server = smtplib.SMTP("smtp.gmail.com",587)
    server.starttls()
    server.login('centre.declaration@gmail.com','yassine@2002')
    msg = EmailMessage()
    message=f"\n code est : {randnbr}"
    msg.set_content(message)
    msg['Subject']="code de verification :"
    msg['From']='centre.declaration@gmail.com'
    msg['To']=email
    server.sendmail('centre.declaration@gmail.com',email,msg.as_string())
    server.quit()



def add_new_user(super_usr,pw,secre,email,key_id,ten_id):
    config = configparser.ConfigParser()
    cc=email.split('@')[0]
    config.add_section(cc)
    config.set(cc, 'email', email)
    config.set(cc, 'passwd', pw)
    config.set(cc, 'cle_secret', secre)
    config.set(cc, 'cle_id', key_id)
    config.set(cc, 'ten', ten_id)
    with open(os.getcwd()+"\\dashbord\\__pycache__\\configfile.ini", 'w') as configfile:
        config.write(configfile)

    msg = EmailMessage()
    msg['Subject'] = 'New notification to add user'
    msg['From'] = 'centre.declaration@gmail.com' 
    msg['To'] = super_usr

    msg.set_content('''
    <!DOCTYPE html>
    <html>
        <body>
            <div style="background-color:#eee;padding:10px 20px;">
                <h2 style="font-family:Georgia, 'Times New Roman', Times, serif;color#454349;">New Request To Login</h2>
            </div>
            <div style="padding:20px 0px">
                <div style="height: 500px;width:400px">
                    <div style="text-align:center;">
                        <h3>Add New User</h3>
                        <p>We want to declare to you that there's new user want to have access to web</p><br>
                        <a href="http://127.0.0.1:8000/add_user/?q1='''+cc+'''">Add</a>
                    </div>
                </div>
            </div>
        </body>
    </html>
    ''', subtype='html')


    with smtplib.SMTP_SSL('smtp.gmail.com', 465) as smtp:
        smtp.login('centre.declaration@gmail.com','rcrtbuvmjyxkofcj') 
        smtp.send_message(msg)


